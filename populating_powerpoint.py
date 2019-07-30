# -*- coding: utf-8 -*-
"""
Created on Wed Jul 17 14:32:30 2019

@author: payam.bagheri
"""

import pandas as pd
import numpy as np
from os import path
from pptx import Presentation
from pptx.chart.data import CategoryChartData
import pptx
from pptx.util import Inches

dir_path = path.dirname(path.dirname(path.abspath(__file__)))
print(dir_path)
data = pd.read_excel(dir_path + '/0_input_data/2067-love-data.xlsx')
data['know-to-lovelike'] = data['love-like']/data['know']
data['lovelike-to-love'] = data['love']/data['love-like']

prs = Presentation(dir_path + '/0_input_data/2067_brand_card.pptx')

demo_cols = ['Male', 'Female', '18-24', '25-39', 'West', 'Ontario', 'Quebec']
indices = pd.DataFrame(columns = demo_cols, index = range(data.shape[0]))

for col in demo_cols:
    indices[col] = data[col]/data['love']*100


    
#indices.to_csv(dir_path + '/0_output/indices.csv', index = False)


'''
for i in range(len(prs.slides[0].shapes)):
    print('i is ', i, prs.slides[0].shapes[i].name)
    
prs.slides[0].shapes[15].name

for i in range(len(prs.slides[0].shapes)):
    if prs.slides[0].shapes[i].has_table:
        print('i is ', i, prs.slides[0].shapes[i].name)
        #print(prs.slides[0].shapes[i].has_table)


for i in range(len(prs.slides[0].shapes)):
    try:
        print('i is ', i)
        print(prs.slides[0].shapes[i].text)
    except AttributeError:
        pass
    
prs.slides[0].shapes[17].text = 'Payam'

for i in range(len(prs.slides[0].shapes)):
    if prs.slides[0].shapes[i].has_chart:
        print('i is ', i)
        #print(prs.slides[0].shapes[i].has_chart)
        
for i in range(len(prs.slides[0].shapes)):
    if prs.slides[0].shapes[i].has_table:
        print('i is ', i)
        #print(prs.slides[0].shapes[i].has_table)
'''

for i, brand in enumerate(data['brand'][0:]):
    print(i, brand)
    prs.slides[0].shapes[2].text = 'BRAND LOVE SUMMARY – '
    

    text_frame = prs.slides[0].shapes[2].text_frame
    #text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = brand
    font = run.font
    font.italic = True
    run.text = brand
    font = run.font
    font.name = 'Arial'
    font.size = pptx.util.Pt(20)
    font.bold = True
    #font.italic = None  # cause value to be inherited from theme
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    font.color.rgb = pptx.dml.color.RGBColor(195, 161, 2)
    #prs.slides[0].shapes[2].text = 'Descriptive text'
    
       
    text_frame = prs.slides[0].shapes[3].text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = 'Descriptive text'
    font = run.font
    font.name = 'Arial'
    font.size = pptx.util.Pt(18)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    font.color.rgb = pptx.dml.color.RGBColor(0, 0, 0)


    text_frame = prs.slides[0].shapes[6].table.cell(0,1).text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "{:.0%}".format(data['love'][i])
    font = run.font
    font.name = 'Calibri'
    font.size = pptx.util.Pt(20)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    font.color.rgb = pptx.dml.color.RGBColor(0, 0, 0)
    
    text_frame = prs.slides[0].shapes[6].table.cell(1,1).text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = str(data['rank'][i])
    font = run.font
    font.name = 'Calibri'
    font.size = pptx.util.Pt(20)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    font.color.rgb = pptx.dml.color.RGBColor(0, 0, 0)
    
    '''
    text_frame = prs.slides[0].shapes[21].text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = brand
    font = run.font
    run.text = brand
    font.name = 'Arial'
    font.size = pptx.util.Pt(18)
    font.bold = True
    #font.italic = None  # cause value to be inherited from theme
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    font.color.rgb = pptx.dml.color.RGBColor(0, 0, 0)
    #prs.slides[0].shapes[2].text = 'Descriptive text'
    '''
    
    #img_path = dir_path + '/0_input_data/logos/'+ brand + '.png'
    
    #img = prs.slides[0].shapes[30]
    #sp = textbox.element
    #img.getparent().remove(img)
    #img.delete()
    #left = top = Inches(1)
    #pic = prs.slides[0].shapes.add_picture(img_path, left, top)
    
    #left = Inches(5)
    #height = Inches(5.5)
    #pic = prs.slides[0].shapes.add_picture(img_path, left, top, height=height)
    
    

    
    ftl = ['Male', 'Female', '18-24', '25-39', 'West', 'Ontario', 'Quebec']
    for c in range(len(ftl)):
        text_frame = prs.slides[0].shapes[10].table.cell(c+1,1).text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = str((format(int(round(indices[ftl[c]][i])), 'd')))
        font = run.font
        font.name = 'Arial'
        font.size = pptx.util.Pt(15)
        font.italic = None  # cause value to be inherited from theme
        #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
        if int(round(indices[ftl[c]][i])) > 120:
            font.bold = True
            font.color.rgb = pptx.dml.color.RGBColor(0,176,80)
        elif int(round(indices[ftl[c]][i])) < 80:
            font.bold = True
            font.color.rgb = pptx.dml.color.RGBColor(255, 0, 0)
        
        
    chart_data = CategoryChartData()
    chart_data.categories = 'category_name'
    single_df = pd.DataFrame(columns=['col'], index = range(3))
    single_df.col = list(np.array(data[['know', 'love-like', 'love']][data['brand'] == brand])[0])
    
    for col_idx, col in enumerate(single_df.columns):
        chart_data.add_series(col, (single_df.iloc[:, col_idx].values))
    
    prs.slides[0].shapes[14].chart.replace_data(chart_data)
    prs.slides[0].shapes[14].chart.series[0].data_labels.number_format = '0%'
        

        
       
          
  
    
    chart_data = CategoryChartData()
    chart_data.categories = 'category_name'
    single_df = pd.DataFrame(columns=['col'], index = range(4))
    single_df.col = list(np.array(data[['meets needs', 'unique', 'sets trends', 'trustworthy']][data['brand'] == brand])[0])
    
    for col_idx, col in enumerate(single_df.columns):
        chart_data.add_series(col, (single_df.iloc[:, col_idx].values))
    
    prs.slides[0].shapes[21].chart.replace_data(chart_data)
    prs.slides[0].shapes[21].chart.series[0].data_labels.number_format = '0%'
    
    
    for j, k in [(16, 'know-to-lovelike'), (17, 'lovelike-to-love')]:
        text_frame = prs.slides[0].shapes[j].text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        run = p.add_run()
        #print(data[k][data['brand'] == brand])
        run.text = str("{0:.0%}".format(round(float(data[k][data['brand'] == brand]),2)))
        font = run.font
        font.name = 'Calibri'
        font.size = pptx.util.Pt(16)
        font.bold = False
        font.italic = None  # cause value to be inherited from theme
        #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
        font.color.rgb = pptx.dml.color.RGBColor(0, 0, 0)
    
    prs.save(dir_path + '/0_output/' + brand + '.pptx')
